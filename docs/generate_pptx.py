"""Generate docs/presentation.pptx — a PowerPoint version of the HTML deck.

Run from repo root:
    python docs/generate_pptx.py

Produces: docs/presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- palette (matches HTML deck) ----------
BG_DARK      = RGBColor(0x0B, 0x10, 0x20)
BG_SURFACE   = RGBColor(0x11, 0x18, 0x27)
BG_SURFACE_2 = RGBColor(0x0E, 0x16, 0x27)
BORDER       = RGBColor(0x1F, 0x2A, 0x44)
BORDER_STR   = RGBColor(0x2B, 0x3A, 0x5F)
TEXT         = RGBColor(0xE6, 0xED, 0xF7)
TEXT_DIM     = RGBColor(0x9A, 0xA7, 0xBD)
TEXT_FAINT   = RGBColor(0x6B, 0x78, 0x90)
CYAN         = RGBColor(0x22, 0xD3, 0xEE)
EMERALD      = RGBColor(0x34, 0xD3, 0x99)
VIOLET       = RGBColor(0xA7, 0x8B, 0xFA)
AMBER        = RGBColor(0xFB, 0xBF, 0x24)
ROSE         = RGBColor(0xFB, 0x71, 0x85)
SKY          = RGBColor(0x7D, 0xD3, 0xFC)

FONT_SANS = "Calibri"
FONT_MONO = "Consolas"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # fully blank layout

# ---------- low-level helpers ----------

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    return shp

def add_text(slide, x, y, w, h, text, *, size=16, color=TEXT, bold=False,
             italic=False, font=FONT_SANS, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_paragraphs(slide, x, y, w, h, runs_per_line, *, line_space=1.15,
                   anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """runs_per_line: list of lines; each line is a list of (text, opts) tuples.
    opts: size, color, bold, italic, font."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs_per_line):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_space
        for text, opts in line:
            r = p.add_run()
            r.text = text
            r.font.name = opts.get("font", FONT_SANS)
            r.font.size = Pt(opts.get("size", 14))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", TEXT)
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=12, color=TEXT,
                bullet_color=None, line_space=1.2):
    """items: list of strings. Uses a dot prefix (no native bullets to avoid
    template dependencies)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_space
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "\u2022  "
        r1.font.name = FONT_SANS
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color or CYAN
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT_SANS
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb

def add_accent_bar(slide, x, y, h=Inches(0.06), w=Inches(0.6), color=CYAN):
    add_rect(slide, x, y, w, h, fill=color, line=None)

def add_card(slide, x, y, w, h, *, accent=CYAN, fill=BG_SURFACE, left_bar=True):
    card = add_rect(slide, x, y, w, h, fill=fill, line=BORDER, line_w=Emu(6350))
    if left_bar:
        add_rect(slide, x, y, Inches(0.05), h, fill=accent, line=None)
    return card

def add_pill(slide, x, y, label, color=CYAN):
    w = Inches(1.2); h = Inches(0.26)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid(); shp.fill.fore_color.rgb = BG_SURFACE_2
    shp.line.color.rgb = color
    shp.line.width = Emu(6350)
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT_SANS; r.font.size = Pt(9)
    r.font.color.rgb = color; r.font.bold = True
    return shp

def add_pill_row(slide, x, y, pills):
    """pills: list of (label, color). Returns end-x."""
    cx = x
    for label, color in pills:
        add_pill(slide, cx, y, label, color)
        cx += Inches(1.3)
    return cx

def add_header(slide, eyebrow, title, *, pills=None, title_size=30):
    # accent bar + eyebrow line
    y0 = Inches(0.55)
    if pills:
        add_pill_row(slide, Inches(0.6), y0, pills)
        y0 += Inches(0.36)
    add_text(slide, Inches(0.6), y0, Inches(12), Inches(0.3),
             eyebrow.upper(), size=10, color=CYAN, bold=True)
    add_text(slide, Inches(0.6), y0 + Inches(0.3), Inches(12.2), Inches(0.9),
             title, size=title_size, color=TEXT, bold=True)
    # bottom border line under header
    add_rect(slide, Inches(0.6), y0 + Inches(1.1), Inches(12.2), Emu(6350),
             fill=BORDER, line=None)

def add_footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.1), Inches(6), Inches(0.3),
             "ICESHELF  \u00b7  v1.0.0  \u00b7  Executive + Technical Overview",
             size=9, color=TEXT_FAINT, bold=True)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.2), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=9, color=TEXT_FAINT,
             align=PP_ALIGN.RIGHT, bold=True)

# ---------- table helper ----------

def add_simple_table(slide, x, y, w, h, headers, rows, *,
                     header_color=CYAN, body_size=10, header_size=10):
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    # header
    for j, ht in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = BG_SURFACE_2
        tf = cell.text_frame; tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        p.text = ""
        r = p.add_run(); r.text = ht
        r.font.name = FONT_SANS; r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = header_color
    # body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SURFACE if i % 2 else BG_SURFACE_2
            tf = cell.text_frame
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            # value can be str or (text, color, bold)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            p.text = ""
            if isinstance(val, tuple):
                text, color, bold = val
            else:
                text, color, bold = val, TEXT, False
            r = p.add_run(); r.text = text
            r.font.name = FONT_SANS; r.font.size = Pt(body_size); r.font.bold = bold
            r.font.color.rgb = color
    return tbl

# ---------- slides ----------

def slide_title():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    # brand line
    add_text(s, Inches(0.6), Inches(0.5), Inches(6), Inches(0.3),
             "ICESHELF  \u00b7  v1.0.0", size=11, color=CYAN, bold=True)
    add_text(s, Inches(11), Inches(0.5), Inches(2), Inches(0.3),
             "Apache-2.0", size=11, color=TEXT_FAINT, bold=True, align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
             "APACHE ICEBERG FOR AWS DATA MESH", size=12, color=VIOLET, bold=True)
    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(1.4),
             "Onboard heterogeneous producers\nto Iceberg \u2014 in days, not quarters.",
             size=42, color=TEXT, bold=True)
    add_text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.6),
             "A briefing for senior management and senior architects.",
             size=16, color=TEXT_DIM)

    # badges row
    y = Inches(4.6)
    badges = [
        ("~80% faster producer onboarding", CYAN),
        ("Strict multi-region isolation", VIOLET),
        ("6 AWS stacks, one quality bar", EMERALD),
        ("Apache-2.0 \u00b7 no SaaS dependency", AMBER),
    ]
    cx = Inches(0.6)
    for label, color in badges:
        add_pill(s, cx, y, label, color)
        cx += Inches(3.0)

    # repo + release footer block
    add_rect(s, Inches(0.6), Inches(6.2), Inches(12.2), Emu(6350),
             fill=BORDER, line=None)
    add_text(s, Inches(0.6), Inches(6.35), Inches(7), Inches(0.35),
             "github.com/<org>/Iceshelf", size=14, color=CYAN, bold=True)
    add_text(s, Inches(9), Inches(6.35), Inches(4), Inches(0.35),
             "v1.0.0 \u00b7 Apache-2.0", size=14, color=TEXT, bold=True,
             align=PP_ALIGN.RIGHT)
    return s

def slide_problem():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Problem",
               "Every producer reinvents the same Iceberg wheel.",
               pills=[("Leadership", AMBER), ("Architect", VIOLET)])
    y = Inches(2.2)
    cards = [
        ("Heterogeneous stacks", CYAN,
         ["Glue / EMR PySpark, ECS / Lambda Python, ECS / Lambda Java",
          "Six permutations, six subtly different configuration paths",
          "Each producer re-learns the same quirks"]),
        ("Multi-region under strict rules", VIOLET,
         ["No cross-region S3. No MRAPs. Per-region Glue catalogs.",
          "Metadata repointing is non-trivial and error-prone",
          "Most reference material assumes the opposite posture"]),
        ("No shared operational baseline", EMERALD,
         ["Compaction / snapshot expiry / orphan cleanup diverge per team",
          "Security review reopens from zero on every onboarding",
          "Platform teams cannot enforce policy in one place"]),
    ]
    w = Inches(4.0); h = Inches(3.8); gap = Inches(0.2)
    x = Inches(0.6)
    for title, color, items in cards:
        add_card(s, x, y, w, h, accent=color)
        add_text(s, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
                 title, size=18, color=color, bold=True)
        add_bullets(s, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), h - Inches(1.2),
                    items, size=12, bullet_color=color)
        x += w + gap
    return s

def slide_why_now():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Why now",
               "Iceberg is the de-facto open table format. The window is open.",
               pills=[("Leadership", AMBER), ("Architect", VIOLET)])
    y = Inches(2.2); w = Inches(3.0); h = Inches(3.8); gap = Inches(0.12)
    cards = [
        ("AWS native", CYAN,
         ["Glue 5.1 ships Iceberg-native",
          "Athena engine v3 reads Iceberg",
          "EMR 7.x first-class support"]),
        ("Warehouse convergence", VIOLET,
         ["Snowflake Polaris (open)",
          "Databricks acquired Tabular ($1B+)",
          "Dremio, Starburst, Trino all adopted"]),
        ("Proven at scale", EMERALD,
         ["Netflix, Apple, Airbnb in production",
          "Salesforce: ~50 PB on Iceberg",
          "ASF top-level \u2014 no vendor gate"]),
        ("Regulatory pressure", AMBER,
         ["EU data residency expanding",
          "Audit trails for ML training data",
          "Time travel becoming compliance-relevant"]),
    ]
    x = Inches(0.6)
    for title, color, items in cards:
        add_card(s, x, y, w, h, accent=color)
        add_text(s, x + Inches(0.25), y + Inches(0.3), w - Inches(0.5), Inches(0.5),
                 title, size=16, color=color, bold=True)
        add_bullets(s, x + Inches(0.25), y + Inches(1.0), w - Inches(0.5), h - Inches(1.2),
                    items, size=11, bullet_color=color)
        x += w + gap
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
             "Cost of waiting: every quarter of drift is another wave of producers that must be re-onboarded later.",
             size=12, color=AMBER, italic=True)
    return s

def slide_solution():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Solution",
               "A Claude Code plugin \u2014 skills, subagents, presets.",
               pills=[("All audiences", CYAN)])
    y = Inches(2.2); w = Inches(4.0); h = Inches(3.9); gap = Inches(0.2)
    cards = [
        ("Skills", CYAN,
         ["/iceberg-profile-setup",
          "/iceberg-onboard",
          "/iceberg-ddl  /  /iceberg-data",
          "/iceberg-maintenance  /  /iceberg-info",
          "/iceberg-migrate",
          "/iceberg-multi-region",
          "/iceberg-pipeline"]),
        ("Subagents", VIOLET,
         ["iceberg-orchestrator (end-to-end)",
          "iceberg-architect",
          "iceberg-code-generator",
          "iceberg-validator",
          "iceberg-multi-region-planner"]),
        ("Presets & profiles", EMERALD,
         ["aws-glue-datamesh-strict (default)",
          "aws-glue-central-lake",
          "Per-producer profile YAML",
          "Org-custom presets supported",
          "One source of truth per producer"]),
    ]
    x = Inches(0.6)
    for title, color, items in cards:
        add_card(s, x, y, w, h, accent=color)
        add_text(s, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
                 title, size=18, color=color, bold=True)
        add_bullets(s, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), h - Inches(1.2),
                    items, size=12, bullet_color=color, line_space=1.3)
        x += w + gap
    return s

def slide_exec_summary():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Executive summary",
               "Outcomes, not features.",
               pills=[("Leadership", AMBER)])
    # four big stats
    stats = [
        ("~80%", "faster producer onboarding vs DIY", CYAN),
        ("0", "cross-region S3 reads \u00b7 0 MRAPs", VIOLET),
        ("6", "AWS stacks supported first-class", EMERALD),
        ("100%", "generated code reviewable in a PR", AMBER),
    ]
    y = Inches(2.4); w = Inches(3.0); h = Inches(2.8); gap = Inches(0.12)
    x = Inches(0.6)
    for val, lbl, color in stats:
        add_card(s, x, y, w, h, accent=color, left_bar=False)
        add_rect(s, x, y, w, Inches(0.08), fill=color, line=None)
        add_text(s, x + Inches(0.25), y + Inches(0.4), w - Inches(0.5), Inches(1.2),
                 val, size=48, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), y + Inches(1.9), w - Inches(0.5), Inches(0.8),
                 lbl, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)
        x += w + gap
    # tagline
    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
             "Iceshelf is a code generator with opinions \u2014 your code, your account, your catalog.",
             size=16, color=TEXT, italic=True)
    add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             "No runtime service to buy, no SaaS dependency, no data leaves your AWS account.",
             size=13, color=TEXT_DIM)
    return s

def slide_business_case():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Business case",
               "Illustrative economics for 25 producers.",
               pills=[("Leadership", AMBER)])
    # left: assumptions
    add_card(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(4.6), accent=VIOLET)
    add_text(s, Inches(0.85), Inches(2.4), Inches(5.5), Inches(0.5),
             "Assumptions (illustrative)", size=16, color=VIOLET, bold=True)
    add_bullets(s, Inches(0.85), Inches(3.0), Inches(5.5), Inches(3.4), [
        "25 producers across the mesh over 12 months",
        "DIY baseline: ~10\u201312 person-weeks per producer (design, code, security review, multi-region)",
        "Iceshelf: ~2 person-weeks per producer (profile, review, deploy)",
        "Blended loaded cost assumption: $3,000 / person-week",
        "Excludes cloud spend (identical in both paths) and licensing (none in either)",
    ], size=12, bullet_color=VIOLET, line_space=1.35)
    # right: numbers
    add_card(s, Inches(7.0), Inches(2.2), Inches(5.7), Inches(4.6), accent=EMERALD)
    add_text(s, Inches(7.25), Inches(2.4), Inches(5.2), Inches(0.5),
             "Savings model", size=16, color=EMERALD, bold=True)
    rows = [
        ("DIY effort",           "~275 person-weeks",  "~$825K"),
        ("Iceshelf effort",       "~50 person-weeks",  "~$150K"),
        ("Platform-team overhead","Preset + reviews",  "~$75K"),
        ("Net savings",           "",                  "~$600K"),
        ("Payback",               "",                  "< 3 producers"),
    ]
    add_simple_table(s, Inches(7.25), Inches(3.0), Inches(5.2), Inches(3.3),
                     ["Item", "Units", "Cost"], rows,
                     header_color=EMERALD, body_size=11)
    add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.3),
             "Numbers are illustrative and should be recalibrated with your actual producer count and loaded cost.",
             size=10, color=TEXT_FAINT, italic=True)
    return s

def slide_compare():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "How Iceshelf compares",
               "Different problem, different shape.",
               pills=[("Leadership", AMBER), ("Architect", VIOLET)])
    headers = ["Option", "Runtime lock-in", "Multi-region (strict)", "Cost model"]
    rows = [
        [("Iceshelf (this)", CYAN, True), "None \u2014 generates your code", "First-class, per-region catalogs",  "One-time platform build"],
        ["DIY per producer",              "None",                             "Each team reinvents",                "Repeated cost per team"],
        ["Databricks / Unity + Tabular",  "Databricks runtime",               "Limited; region model differs",      "Per-DBU + data egress"],
        ["Snowflake Polaris / Open Cat.", "Snowflake account",                "Supported via replication \u2014 different shape", "Snowflake credits"],
        ["AWS Glue native",               "AWS managed",                      "DIY multi-region",                    "Glue + S3 pay-as-you-go"],
        ["Dremio / Starburst",            "Query-engine licence",             "Supported; managed tier adds cost",   "Per-node / per-query tier"],
    ]
    add_simple_table(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(3.8),
                     headers, rows, header_color=CYAN, body_size=11, header_size=12)
    # footer cards
    y = Inches(6.15); w = Inches(4.0); gap = Inches(0.1); x = Inches(0.6)
    notes = [
        ("Complement, not replace", EMERALD,
         "Keep your query engines \u2014 Iceshelf only handles producer side."),
        ("Zero egress", CYAN,
         "All data stays in your S3. No vendor SaaS ingress."),
        ("Exit cost: zero", AMBER,
         "Artifacts are your code in your repo. Stop using Iceshelf \u2014 keep the pipelines."),
    ]
    for title, color, body in notes:
        add_card(s, x, y, w, Inches(1.2), accent=color, left_bar=True)
        add_text(s, x + Inches(0.25), y + Inches(0.1), w - Inches(0.5), Inches(0.35),
                 title, size=12, color=color, bold=True)
        add_text(s, x + Inches(0.25), y + Inches(0.5), w - Inches(0.5), Inches(0.7),
                 body, size=10, color=TEXT_DIM)
        x += w + gap
    return s

def slide_audiences():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Who it's for",
               "Three audiences, one platform.",
               pills=[("All audiences", CYAN)])
    y = Inches(2.2); w = Inches(4.0); h = Inches(3.9); gap = Inches(0.2)
    cards = [
        ("Leadership", AMBER,
         ["Faster, safer producer onboarding",
          "One security / audit baseline across the mesh",
          "Predictable unit economics on Iceberg adoption",
          "No runtime lock-in; portable artifacts"]),
        ("Senior Architects", VIOLET,
         ["Strict multi-region (no cross-region S3, no MRAPs)",
          "Preset inheritance for policy-as-code",
          "Glue catalog as source of truth; S3FileIO; KMS",
          "Deterministic metadata repointing utility"]),
        ("Producer developers", EMERALD,
         ["Six stacks supported \u2014 PySpark / PyIceberg / Java",
          "Generated code, not opaque abstractions",
          "Recipes per stack with working examples",
          "Slash commands for day-2 operations"]),
    ]
    x = Inches(0.6)
    for title, color, items in cards:
        add_card(s, x, y, w, h, accent=color)
        add_text(s, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
                 title, size=18, color=color, bold=True)
        add_bullets(s, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), h - Inches(1.2),
                    items, size=12, bullet_color=color, line_space=1.35)
        x += w + gap
    return s

def slide_scope():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Scope v1.0",
               "What's in the box.",
               pills=[("Architect", VIOLET)])
    y = Inches(2.2); w = Inches(6.0); h = Inches(4.5); gap = Inches(0.2)
    in_items = [
        "AWS Glue Data Catalog as Iceberg catalog",
        "Format-version 2 by default (row-level deletes, equality deletes)",
        "Six producer stacks (Glue/EMR PySpark, ECS/Lambda Python/Java)",
        "Strict multi-region: CRR + metadata repointing + per-region catalogs",
        "Presets: aws-glue-datamesh-strict, aws-glue-central-lake",
        "Compaction, snapshot expiry, orphan cleanup generated per producer",
        "IAM / KMS / VPC endpoints encoded in generated Terraform",
    ]
    out_items = [
        "Non-AWS catalogs (Nessie, Hive Metastore, REST) \u2014 future",
        "Active-active multi-region writes \u2014 out of scope by design",
        "Consumer-side query engines \u2014 use what you already have",
        "Observability UI \u2014 generates CloudWatch dashboards, no SaaS UI",
        "Tokenization / masking \u2014 assumed upstream in the pipeline",
    ]
    add_card(s, Inches(0.6), y, w, h, accent=CYAN)
    add_text(s, Inches(0.85), y + Inches(0.3), w - Inches(0.5), Inches(0.5),
             "In scope", size=18, color=CYAN, bold=True)
    add_bullets(s, Inches(0.85), y + Inches(1.0), w - Inches(0.5), h - Inches(1.2),
                in_items, size=12, bullet_color=CYAN, line_space=1.3)
    add_card(s, Inches(0.6) + w + gap, y, w, h, accent=AMBER)
    add_text(s, Inches(0.85) + w + gap, y + Inches(0.3), w - Inches(0.5), Inches(0.5),
             "Out of scope (today)", size=18, color=AMBER, bold=True)
    add_bullets(s, Inches(0.85) + w + gap, y + Inches(1.0), w - Inches(0.5), h - Inches(1.2),
                out_items, size=12, bullet_color=AMBER, line_space=1.3)
    return s

def slide_architecture():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Architecture",
               "One producer, one region \u2014 the base shape.",
               pills=[("Architect", VIOLET), ("Developer", EMERALD)])
    # horizontal flow: Claude Code -> Iceshelf -> Generated code -> AWS
    y = Inches(2.4); h = Inches(3.0)
    boxes = [
        ("Claude Code", "Interactive CLI\nSkills + subagents", CYAN),
        ("Iceshelf", "Profile-first generation\nPresets + templates", VIOLET),
        ("Generated code", "PySpark / PyIceberg / Java\nTerraform / runbooks", EMERALD),
        ("Producer AWS", "Glue catalog \u00b7 S3 warehouse\nCloudWatch \u00b7 KMS \u00b7 VPC", AMBER),
    ]
    w = Inches(2.85); gap = Inches(0.28)
    x = Inches(0.6)
    for i, (title, body, color) in enumerate(boxes):
        add_card(s, x, y, w, h, accent=color)
        add_text(s, x + Inches(0.25), y + Inches(0.3), w - Inches(0.5), Inches(0.5),
                 title, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), y + Inches(1.2), w - Inches(0.5), Inches(1.8),
                 body, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            # arrow
            ax = x + w + Emu(10000)
            ay = y + h / 2
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        ax, ay - Inches(0.15),
                                        gap - Emu(20000), Inches(0.3))
            arrow.shadow.inherit = False
            arrow.fill.solid(); arrow.fill.fore_color.rgb = BORDER_STR
            arrow.line.fill.background()
        x += w + gap
    # bottom note
    add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
             "Iceshelf runs at authoring time. Nothing from Iceshelf sits on the runtime data path.",
             size=13, color=TEXT, italic=True)
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "Producer code talks directly to Glue and S3 \u2014 no proxy, no sidecar, no external service.",
             size=12, color=TEXT_DIM)
    return s

def slide_multi_region():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Multi-region \u00b7 the differentiator",
               "DR resilience without cross-region S3 or MRAPs.",
               pills=[("Architect", VIOLET), ("Leadership", AMBER)])
    # two region cards + arrow
    y = Inches(2.2); h = Inches(3.6)
    pw = Inches(5.4)
    # primary
    add_card(s, Inches(0.6), y, pw, h, accent=CYAN)
    add_text(s, Inches(0.85), y + Inches(0.25), pw - Inches(0.5), Inches(0.4),
             "PRIMARY \u00b7 us-east-1", size=12, color=CYAN, bold=True)
    add_bullets(s, Inches(0.85), y + Inches(0.8), pw - Inches(0.5), h - Inches(1.0), [
        "S3 warehouse bucket \u2014 data + metadata written locally",
        "Glue catalog \u2014 table registered with primary metadata_location",
        "Ingestion job \u2014 glue.region + s3.region pinned locally",
    ], size=12, bullet_color=CYAN, line_space=1.35)
    # DR
    add_card(s, Inches(7.3), y, pw, h, accent=VIOLET)
    add_text(s, Inches(7.55), y + Inches(0.25), pw - Inches(0.5), Inches(0.4),
             "DR \u00b7 us-west-2", size=12, color=VIOLET, bold=True)
    add_bullets(s, Inches(7.55), y + Inches(0.8), pw - Inches(0.5), h - Inches(1.0), [
        "Replicated bucket \u2014 same keyspace via S3 CRR",
        "Repointing utility \u2014 rewrites metadata.json + manifests",
        "Local Glue catalog \u2014 points at repointed metadata",
        "Queryable in-region, zero cross-region I/O at runtime",
    ], size=12, bullet_color=VIOLET, line_space=1.35)
    # CRR arrow in the middle
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(6.05), Inches(3.6),
                               Inches(1.2), Inches(0.4))
    arrow.shadow.inherit = False
    arrow.fill.solid(); arrow.fill.fore_color.rgb = EMERALD
    arrow.line.fill.background()
    add_text(s, Inches(6.05), Inches(3.25), Inches(1.2), Inches(0.3),
             "S3 CRR", size=10, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.05), Inches(4.05), Inches(1.2), Inches(0.3),
             "repoint \u2190", size=10, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

    # footer tags
    tags = [
        ("Zero cross-region I/O", CYAN),
        ("No MRAPs", VIOLET),
        ("Per-region queryable", EMERALD),
        ("Auditable failover", AMBER),
    ]
    y2 = Inches(6.1); x = Inches(0.6); w = Inches(3.0); gap = Inches(0.1)
    for label, color in tags:
        add_card(s, x, y2, w, Inches(0.8), accent=color, left_bar=True)
        add_text(s, x + Inches(0.25), y2 + Inches(0.22), w - Inches(0.5), Inches(0.4),
                 label, size=12, color=color, bold=True)
        x += w + gap
    return s

def slide_rpo_rto():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "RPO / RTO & failure modes",
               "Numbers you can put on a risk register.",
               pills=[("Architect", VIOLET), ("Leadership", AMBER)])
    # two tables
    y = Inches(2.15); h = Inches(4.2)
    # left: recovery targets
    add_text(s, Inches(0.6), y, Inches(6.0), Inches(0.4),
             "Recovery targets (default profile)", size=14, color=CYAN, bold=True)
    targets = [
        ["RPO \u2014 data loss window",          ("\u2264 20 min", EMERALD, True),  "CRR lag p99 + last repointed snapshot"],
        ["RTO \u2014 time to query in DR",       ("\u2264 60 min", EMERALD, True),  "Repointing + Glue re-register"],
        ["Repointing cadence",                   "every 15 min",                      "EventBridge \u2192 Lambda / SFN"],
        ["CRR SLA (AWS)",                        "15 min / 99%",                      "S3 RTC, opt-in per profile"],
        ["Snapshot retention (DR)",              "7 days default",                    "Override in preset"],
        ["Failover decision SLA",                "\u2264 15 min",                     "Human gate; alarms page on-call"],
    ]
    add_simple_table(s, Inches(0.6), Inches(2.5), Inches(6.1), Inches(3.7),
                     ["Metric", "Target", "Driver"], targets,
                     header_color=CYAN, body_size=10)
    # right: failure modes
    add_text(s, Inches(7.0), y, Inches(6.0), Inches(0.4),
             "Failure-mode matrix", size=14, color=VIOLET, bold=True)
    fm = [
        ["CRR lag spike (>30 min)",              "BytesPendingReplication alarm",  "Skip commit; DR stays last-known-good"],
        ["Repointing fails mid-run",             "Step Functions task failure",     "Atomic swap; prior metadata retained"],
        ["Primary writes during failover",       "Runbook write-disable gate",      "IAM deny \u2192 final repoint \u2192 DR"],
        ["Stale DR catalog at query time",       "Snapshot-age metric",             "Last repointed snapshot served"],
        ["Orphaned data in DR bucket",           "Monthly orphan-cleanup job",      "GC older than retention"],
        ["Failback to primary",                  "Reverse CRR + repoint back",      "Same utility, flipped direction"],
    ]
    add_simple_table(s, Inches(7.0), Inches(2.5), Inches(5.7), Inches(3.7),
                     ["Scenario", "Detection", "Behavior"], fm,
                     header_color=VIOLET, body_size=10)
    add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.3),
             "Targets are knobs, not guesses \u2014 CloudWatch alarms generated to prove them. Active-active is out of scope by design.",
             size=10, color=TEXT_FAINT, italic=True)
    return s

def slide_security():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Security, compliance & governance",
               "Defensible by default \u2014 built for security review.",
               pills=[("Architect", VIOLET), ("Leadership", AMBER)])
    y = Inches(2.2); w = Inches(6.0); h = Inches(2.3); gap = Inches(0.12)
    cards = [
        ("Identity & access", CYAN,
         ["IAM scoped to specific DB + table ARNs (never *)",
          "Lake Formation fine-grained via profile (lf_tags)",
          "Separate read / write / maintenance roles per producer",
          "No cross-account catalog writes by default"]),
        ("Encryption & key management", VIOLET,
         ["S3 bucket encryption defaults to aws:kms (CMK alias pinned)",
          "Glue catalog + connection-credential encryption supported",
          "Per-producer CMK isolation \u2014 a one-line profile change",
          "Key-grant list in profile \u2014 rotation / audit scripted"]),
        ("Network isolation", EMERALD,
         ["VPC endpoints for S3 / Glue / STS / KMS / CloudWatch",
          "No public egress to any Iceshelf-hosted service (none exists)",
          "Private-subnet-only deployments out of the box",
          "Multi-region data path stays in-region; CRR is the only hop"]),
        ("Audit & governance", AMBER,
         ["Profile + artifacts diff-reviewable in Git \u2014 every change attributable",
          "Required tags enforced at generation (CostCenter, Owner, \u2026)",
          "Optional audit bucket: WORM mirror of every commit event",
          "Snapshot expiry + orphan cleanup with retention windows"]),
    ]
    positions = [(Inches(0.6), y), (Inches(0.6) + w + gap, y),
                 (Inches(0.6), y + h + gap), (Inches(0.6) + w + gap, y + h + gap)]
    for (x, yy), (title, color, items) in zip(positions, cards):
        add_card(s, x, yy, w, h, accent=color)
        add_text(s, x + Inches(0.25), yy + Inches(0.2), w - Inches(0.5), Inches(0.4),
                 title, size=14, color=color, bold=True)
        add_bullets(s, x + Inches(0.25), yy + Inches(0.7), w - Inches(0.5), h - Inches(0.8),
                    items, size=10, bullet_color=color, line_space=1.25)
    # framework strip
    fw_y = y + 2 * h + 2 * gap + Inches(0.1)
    add_card(s, Inches(0.6), fw_y, Inches(12.1), Inches(0.9), accent=CYAN, left_bar=False)
    fws = [("SOC 2", CYAN), ("HIPAA / PHI", VIOLET),
           ("PCI-DSS", EMERALD), ("FedRAMP / GovCloud", AMBER)]
    x = Inches(0.85); cell_w = Inches(2.95)
    for label, color in fws:
        add_text(s, x, fw_y + Inches(0.15), cell_w, Inches(0.3),
                 label, size=12, color=color, bold=True)
        add_text(s, x, fw_y + Inches(0.45), cell_w, Inches(0.4),
                 "aligned by design", size=10, color=TEXT_DIM)
        x += cell_w + Inches(0.05)
    return s

def slide_tech_stacks():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Tech stack coverage",
               "Six producer stacks. One quality bar.",
               pills=[("Developer", EMERALD)])
    y = Inches(2.2); w = Inches(4.0); h = Inches(2.3); gap = Inches(0.12)
    cards = [
        ("Glue PySpark", CYAN,   "AWS Glue 5.0 / 5.1\nSpark 3.5.x\niceberg-spark-runtime 1.10"),
        ("EMR PySpark",  CYAN,   "EMR 6.5+ / 7.x\nScala 2.12 / 2.13\niceberg-spark-runtime 1.10"),
        ("ECS Python",   EMERALD,"Fargate long-running\npyiceberg[glue,s3fs,pyarrow]\u2265 0.8"),
        ("Lambda Python",EMERALD,"Lambda \u2264 15 min\nReserved concurrency = 1\npyiceberg 0.8+"),
        ("ECS Java",     VIOLET, "JVM long-running\niceberg-core + iceberg-aws 1.10\nGlueCatalog + S3FileIO"),
        ("Lambda Java",  VIOLET, "JVM 21 runtime\nSnapStart for cold-start\nSingle-writer pattern"),
    ]
    for i, (title, color, body) in enumerate(cards):
        row, col = divmod(i, 3)
        x = Inches(0.6) + col * (w + gap)
        yy = y + row * (h + gap)
        add_card(s, x, yy, w, h, accent=color)
        add_text(s, x + Inches(0.3), yy + Inches(0.25), w - Inches(0.5), Inches(0.5),
                 title, size=16, color=color, bold=True)
        add_text(s, x + Inches(0.3), yy + Inches(0.85), w - Inches(0.5), h - Inches(1.0),
                 body, size=12, color=TEXT_DIM, font=FONT_MONO)
    return s

def slide_skills():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Skills",
               "Slash commands, one per job.",
               pills=[("Developer", EMERALD)])
    rows = [
        [("/iceberg-profile-setup", CYAN, True), "Run first \u2014 captures producer conventions into profile YAML."],
        [("/iceberg-onboard",        CYAN, True), "End-to-end onboarding assessment for new or existing producers."],
        [("/iceberg-ddl",            CYAN, True), "Create, alter, drop Iceberg tables with schema evolution."],
        [("/iceberg-data",           CYAN, True), "Insert, upsert, delete data with MoR / CoW guidance."],
        [("/iceberg-maintenance",    CYAN, True), "Compaction, snapshot expiry, orphan cleanup."],
        [("/iceberg-info",           CYAN, True), "Table metadata inspection and diagnostics."],
        [("/iceberg-migrate",        CYAN, True), "Migrate existing Parquet / Hive tables to Iceberg."],
        [("/iceberg-multi-region",   CYAN, True), "Multi-region resilience setup with metadata repointing."],
        [("/iceberg-pipeline",       CYAN, True), "End-to-end pipeline generation for new or existing producers."],
    ]
    add_simple_table(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(4.6),
                     ["Skill", "Purpose"], rows, header_color=CYAN, body_size=12)
    return s

def slide_subagents():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Subagents",
               "Specialized, coordinated.",
               pills=[("Architect", VIOLET)])
    rows = [
        [("iceberg-orchestrator",        VIOLET, True), "Coordinates end-to-end onboarding. Delegates to the specialists below."],
        [("iceberg-architect",            VIOLET, True), "Analyzes producer setup and designs migration strategy."],
        [("iceberg-code-generator",       VIOLET, True), "Generates platform-specific production code and IaC."],
        [("iceberg-validator",            VIOLET, True), "Validates configurations, schemas, and generated code."],
        [("iceberg-multi-region-planner", VIOLET, True), "Plans and generates multi-region infrastructure and sync."],
    ]
    add_simple_table(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(3.0),
                     ["Subagent", "Purpose"], rows, header_color=VIOLET, body_size=12)
    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
             "All subagents share the producer profile \u2014 one source of truth, no re-asking.",
             size=13, color=TEXT, italic=True)
    return s

def slide_producer_flow():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Producer journey",
               "From zero to production in seven phases.",
               pills=[("Developer", EMERALD), ("Architect", VIOLET)])
    phases = [
        ("1", "Profile setup",      "Capture producer conventions"),
        ("2", "Assessment",         "Onboard scan \u2014 existing state"),
        ("3", "DDL design",         "Schema + partitioning + format-v2"),
        ("4", "Data path",          "Insert / upsert / delete patterns"),
        ("5", "Maintenance",        "Compaction + snapshot expiry"),
        ("6", "Multi-region",       "CRR + repoint + DR catalog"),
        ("7", "Promote to prod",    "Runbook + alarms + handoff"),
    ]
    y = Inches(2.4); w = Inches(1.77); h = Inches(2.8); gap = Inches(0.06)
    x = Inches(0.6)
    for i, (num, title, body) in enumerate(phases):
        color = [CYAN, VIOLET, EMERALD, AMBER, CYAN, VIOLET, EMERALD][i]
        add_card(s, x, y, w, h, accent=color)
        add_text(s, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.7),
                 num, size=32, color=color, bold=True)
        add_text(s, x + Inches(0.2), y + Inches(1.0), w - Inches(0.4), Inches(0.5),
                 title, size=12, color=TEXT, bold=True)
        add_text(s, x + Inches(0.2), y + Inches(1.6), w - Inches(0.4), Inches(1.2),
                 body, size=10, color=TEXT_DIM)
        x += w + gap
    add_text(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.5),
             "Each phase produces reviewable artifacts \u2014 code, IaC, and runbooks \u2014 committed to the producer's repo.",
             size=12, color=TEXT_DIM)
    return s

def slide_profiles():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Profiles",
               "One YAML per producer. Preset-inherited.",
               pills=[("Architect", VIOLET)])
    # left: sample YAML
    add_card(s, Inches(0.6), Inches(2.2), Inches(7.0), Inches(4.5), accent=EMERALD, left_bar=False)
    add_rect(s, Inches(0.6), Inches(2.2), Inches(7.0), Inches(0.08), fill=EMERALD, line=None)
    add_text(s, Inches(0.85), Inches(2.35), Inches(6.5), Inches(0.4),
             "profiles/datamesh-producer-brok.yaml", size=11, color=EMERALD, bold=True,
             font=FONT_MONO)
    yaml_lines = [
        ("producer_name", "datamesh-producer-brok"),
        ("extends",       "aws-glue-datamesh-strict"),
        ("primary_region","us-east-1"),
        ("dr_region",     "us-west-2"),
        ("warehouse_bucket_primary", "brok-confirmed-us-east-1"),
        ("warehouse_bucket_dr",      "brok-confirmed-us-west-2"),
        ("database",      "brok_confirmed"),
        ("tech_stack",    "glue_pyspark"),
        ("orchestrator",  "step_functions"),
    ]
    lines = [[("# " + "Minimal profile \u2014 only overrides vs the preset", {"size":11,"color":TEXT_FAINT,"italic":True,"font":FONT_MONO})]]
    for k, v in yaml_lines:
        lines.append([
            (k + ": ",    {"size":12,"color":CYAN,"bold":True,"font":FONT_MONO}),
            (v,           {"size":12,"color":TEXT,"font":FONT_MONO}),
        ])
    add_paragraphs(s, Inches(0.85), Inches(2.8), Inches(6.5), Inches(3.8),
                   lines, line_space=1.25)
    # right: why profiles
    add_card(s, Inches(7.8), Inches(2.2), Inches(4.9), Inches(4.5), accent=VIOLET)
    add_text(s, Inches(8.05), Inches(2.4), Inches(4.5), Inches(0.4),
             "Why this shape", size=14, color=VIOLET, bold=True)
    add_bullets(s, Inches(8.05), Inches(3.0), Inches(4.5), Inches(3.7), [
        "Extends a preset \u2014 only overrides appear",
        "Every skill reads from the profile (answers captured once)",
        "Org can publish its own preset; producers inherit it",
        "Diff-reviewable: policy drift is a PR, not a Slack thread",
        "Minimal surface: ~10 fields for the common case",
    ], size=12, bullet_color=VIOLET, line_space=1.35)
    return s

def slide_release_readiness():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Release readiness",
               "v1.0.0 \u2014 what's proven.",
               pills=[("Leadership", AMBER), ("Architect", VIOLET)])
    rows = [
        [("Skills",                CYAN, True), "9 / 9 \u2014 all covered with recipes", ("Ready", EMERALD, True)],
        [("Subagents",             CYAN, True), "5 / 5 \u2014 orchestrator + 4 specialists", ("Ready", EMERALD, True)],
        [("Presets",               CYAN, True), "2 built-in (strict + central-lake) + org-custom hook", ("Ready", EMERALD, True)],
        [("Tech stacks",           CYAN, True), "6 / 6 supported with working code templates", ("Ready", EMERALD, True)],
        [("Multi-region",          CYAN, True), "CRR + repointing + per-region catalogs (strict posture)", ("Ready", EMERALD, True)],
        [("Observability",         CYAN, True), "CloudWatch alarms + dashboards generated per producer", ("Ready", EMERALD, True)],
        [("Documentation",         CYAN, True), "Admin guide + 3 stack recipes + end-user guide", ("Ready", EMERALD, True)],
        [("Security review pkg",   CYAN, True), "IAM / KMS / VPC / audit artifacts generated",          ("Ready", EMERALD, True)],
    ]
    add_simple_table(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(4.4),
                     ["Area", "What ships in v1.0", "Status"], rows,
                     header_color=CYAN, body_size=11)
    add_text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.4),
             "License: Apache-2.0  \u00b7  Artifacts: Claude Code plugin + profile YAML + generated code in your repo.",
             size=11, color=TEXT_FAINT, italic=True)
    return s

def slide_roadmap():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Roadmap",
               "What we're sequencing next.",
               pills=[("Architect", VIOLET), ("Leadership", AMBER)])
    quarters = [
        ("Near-term (next quarter)", CYAN, [
            "Nessie / REST catalog adapters",
            "Glue Data Quality integration",
            "CDC ingestion templates (DMS / DynamoDB streams)",
            "Cost & storage dashboards (per-producer + fleet)",
        ]),
        ("Mid-term (2 quarters)", VIOLET, [
            "Lake Formation row/column-level access generator",
            "Schema-evolution validator with consumer-contract checks",
            "GovCloud profile preset",
            "Producer marketplace for shared DDL packs",
        ]),
        ("Longer-term", AMBER, [
            "Non-AWS cloud targets (GCP / Azure catalogs)",
            "Active-active via conflict-aware partitioning (research)",
            "Query-side recipes: Athena / Trino / Snowflake Polaris federation",
            "Attached notebooks for ad-hoc exploration",
        ]),
    ]
    y = Inches(2.2); w = Inches(4.0); h = Inches(4.5); gap = Inches(0.2)
    x = Inches(0.6)
    for title, color, items in quarters:
        add_card(s, x, y, w, h, accent=color)
        add_text(s, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
                 title, size=16, color=color, bold=True)
        add_bullets(s, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), h - Inches(1.2),
                    items, size=12, bullet_color=color, line_space=1.3)
        x += w + gap
    return s

def slide_pilot():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "The ask \u00b7 90-day pilot",
               "Two producers. One quarter. Clear decision gate at week 12.",
               pills=[("Leadership", AMBER), ("Architect", VIOLET)])
    # left: scope
    add_card(s, Inches(0.6), Inches(2.2), Inches(6.1), Inches(2.4), accent=CYAN)
    add_text(s, Inches(0.85), Inches(2.35), Inches(5.6), Inches(0.4),
             "Scope & shape", size=14, color=CYAN, bold=True)
    add_bullets(s, Inches(0.85), Inches(2.8), Inches(5.6), Inches(1.7), [
        "Two producers \u2014 one PySpark, one PyIceberg",
        "One region pair \u2014 e.g. us-east-1 / us-west-2",
        "One preset \u2014 aws-glue-datamesh-strict by default",
        "Dry-run first \u2014 code reviewed in a PR before prod",
    ], size=11, bullet_color=CYAN, line_space=1.3)
    # right: timeline
    add_card(s, Inches(6.9), Inches(2.2), Inches(5.8), Inches(2.4), accent=VIOLET)
    add_text(s, Inches(7.15), Inches(2.35), Inches(5.3), Inches(0.4),
             "Timeline (12 weeks)", size=14, color=VIOLET, bold=True)
    tl = [
        ["1\u20132",  "Preset sign-off; sandbox; profiles captured"],
        ["3\u20135",  "Producer #1: onboard \u2192 DDL \u2192 ingestion \u2192 maintenance"],
        ["6\u20137",  "Multi-region for Producer #1: CRR + repoint + DR"],
        ["8\u20139",  "Producer #2 on the same preset (reuse proof)"],
        ["10",         "Game-day: primary outage; RPO / RTO on stopwatch"],
        ["11",         "Security, cost, auditor review"],
        ["12",         "Decision gate: go / no-go on mesh-wide rollout"],
    ]
    add_simple_table(s, Inches(7.15), Inches(2.75), Inches(5.3), Inches(1.8),
                     ["Week", "Milestone"], tl, header_color=VIOLET, body_size=9)
    # bottom row: exit criteria + asks
    add_card(s, Inches(0.6), Inches(4.75), Inches(6.1), Inches(2.2), accent=EMERALD)
    add_text(s, Inches(0.85), Inches(4.9), Inches(5.6), Inches(0.4),
             "Exit criteria (all must pass)", size=14, color=EMERALD, bold=True)
    add_bullets(s, Inches(0.85), Inches(5.35), Inches(5.6), Inches(1.55), [
        "Both producers in production with zero hand-edits",
        "\u2264 2 person-weeks per producer (vs 10\u201312 DIY)",
        "RPO \u2264 20 min and RTO \u2264 60 min on game-day",
        "Security / cost reviews signed off",
    ], size=10, bullet_color=EMERALD, line_space=1.3)
    add_card(s, Inches(6.9), Inches(4.75), Inches(5.8), Inches(2.2), accent=AMBER)
    add_text(s, Inches(7.15), Inches(4.9), Inches(5.3), Inches(0.4),
             "What we need from you", size=14, color=AMBER, bold=True)
    add_bullets(s, Inches(7.15), Inches(5.35), Inches(5.3), Inches(1.55), [
        "Executive sponsor \u2014 clears org-level blockers",
        "Platform-team lead \u2014 0.5 FTE for 12 weeks",
        "Two producer teams \u2014 ~0.25 FTE each",
        "Sandbox AWS account with target region pair",
        "Security reviewer for a single walkthrough session",
    ], size=10, bullet_color=AMBER, line_space=1.3)
    return s

def slide_next_steps():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    add_header(s, "Next steps",
               "One action for each audience.",
               pills=[("All audiences", CYAN)])
    y = Inches(2.2); w = Inches(4.0); h = Inches(4.2); gap = Inches(0.2)
    cards = [
        ("Leadership", AMBER,
         ["Approve Iceberg as the mesh-wide table format",
          "Sponsor a 2-producer pilot in the next quarter",
          "Designate the platform-team owner for preset policy"],
         "README.md \u00b7 CHANGELOG.md"),
        ("Senior Architects", VIOLET,
         ["Review docs/admin-guide.md",
          "Pick aws-glue-datamesh-strict or aws-glue-central-lake",
          "Define gating and wave-rollout policy"],
         "profiles/schema.md \u00b7 profiles/presets/*.yaml"),
        ("Producer developers", EMERALD,
         ["Read the recipe for your stack",
          "Run /iceberg-profile-setup in your repo",
          "Run /iceberg-onboard and review the deployment package"],
         "docs/examples/<stack>-producer-recipe.md"),
    ]
    x = Inches(0.6)
    for title, color, items, read in cards:
        add_card(s, x, y, w, h, accent=color)
        add_pill(s, x + Inches(0.3), y + Inches(0.3), title, color)
        add_text(s, x + Inches(0.3), y + Inches(0.75), w - Inches(0.6), Inches(0.5),
                 "What to do next", size=14, color=TEXT, bold=True)
        add_bullets(s, x + Inches(0.3), y + Inches(1.25), w - Inches(0.6), h - Inches(2.2),
                    items, size=12, bullet_color=color, line_space=1.35)
        add_rect(s, x + Inches(0.3), y + h - Inches(0.9), w - Inches(0.6), Emu(6350),
                 fill=BORDER, line=None)
        add_text(s, x + Inches(0.3), y + h - Inches(0.75), w - Inches(0.6), Inches(0.4),
                 "Read:", size=9, color=TEXT_FAINT, bold=True)
        add_text(s, x + Inches(0.3), y + h - Inches(0.5), w - Inches(0.6), Inches(0.4),
                 read, size=10, color=CYAN, font=FONT_MONO)
        x += w + gap
    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "github.com/<org>/Iceshelf  \u00b7  v1.0.0  \u00b7  Apache-2.0",
             size=12, color=CYAN, bold=True)
    return s

# ---------- build ----------

builders = [
    slide_title,
    slide_problem,
    slide_why_now,
    slide_solution,
    slide_exec_summary,
    slide_business_case,
    slide_compare,
    slide_audiences,
    slide_scope,
    slide_architecture,
    slide_multi_region,
    slide_rpo_rto,
    slide_security,
    slide_tech_stacks,
    slide_skills,
    slide_subagents,
    slide_producer_flow,
    slide_profiles,
    slide_release_readiness,
    slide_roadmap,
    slide_pilot,
    slide_next_steps,
]

total = len(builders)
slides_built = []
for fn in builders:
    slides_built.append(fn())

# add footers with index/total now that total is known
for i, sl in enumerate(slides_built, start=1):
    add_footer(sl, i, total)

out = Path(__file__).parent / "presentation.pptx"
prs.save(out)
print(f"Wrote {out}  ({total} slides)")
